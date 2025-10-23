#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os

# Presentation generator for MagdyDiner final project
# Saves to presentation/MagdyDiner_Presentation.pptx

slides = [
    {
        "title": "MagdyDiner - Android App",
        "subtitle": "Final Project Presentation\nPresenter: Magdy (Magdy1996)",
        "notes": "Introduce yourself, the repo and branch (main). Mention this presentation is a short 5-8 minute overview with a short demo."
    },
    {
        "title": "Problem & Goals",
        "bullets": [
            "Problem statement: the user pain the app solves",
            "Primary goals: usability, reliability, offline support",
            "Success criteria: stable build, complete main flows tested"
        ],
        "notes": "Concise problem statement and measurable goals. Keep this slide to 20-30s."
    },
    {
        "title": "Tech Stack & Tools",
        "bullets": [
            "Languages: Kotlin (primary), Java (if present)",
            "Build: Gradle; IDE: Android Studio",
            "Key libraries: Retrofit, Coroutines, Room, Hilt, Jetpack components"
        ],
        "notes": "Call out architecture pattern (MVVM) and why these libraries were chosen."
    },
    {
        "title": "High-level Architecture",
        "bullets": [
            "UI (Activities/Fragments) -> ViewModel -> Repository -> Data sources",
            "Remote API + Local DB (Room) with Repository merging data",
            "Single source of truth; ViewModel exposes StateFlow/LiveData"
        ],
        "notes": "Explain flow of data and where business logic lives. Use a quick diagram while speaking."
    },
    {
        "title": "Key Components & Responsibilities",
        "bullets": [
            "UI: render state, handle user input",
            "ViewModel: UI state, orchestrates calls to repositories",
            "Repository: abstracts data sources",
            "Data sources: API client (Retrofit), local DB (Room)"
        ],
        "notes": "Map important packages/classes and a short responsibility sentence for each."
    },
    {
        "title": "Resource Wrapper (Loading/Success/Error)",
        "bullets": [
            "Purpose: unify network/db result handling",
            "Typical shape: Resource.Success(data) / Resource.Error(msg) / Resource.Loading",
            "UI observes Resource and shows progress / content / error"
        ],
        "notes": "If you have a Resource.kt file, mention its package and show a short snippet during demo."
    },
    {
        "title": "Demo Setup",
        "bullets": [
            "Ensure emulator or physical device connected",
            "Build: ./gradlew assembleDebug",
            "Run from Android Studio or install APK from build/outputs/apk"
        ],
        "notes": "Explain how you will run the demo and what to look for (Logcat filter by package)."
    },
    {
        "title": "Live Demo Script",
        "bullets": [
            "1) Launch app and show main screen",
            "2) Perform key flows (login / search / add favorite / offline behavior)",
            "3) Trigger an error to show error handling and Resource.Error"
        ],
        "notes": "Have emulator prepared, breakpoints set, network inspector available. Use a short, rehearsed script."
    },
    {
        "title": "Tests & CI",
        "bullets": [
            "Unit tests: ./gradlew test",
            "Instrumentation tests: ./gradlew connectedAndroidTest",
            "(Optional) CI: mention GitHub Actions or other if present"
        ],
        "notes": "If you have tests, mention coverage and how to run them quickly."
    },
    {
        "title": "Known Issues & Fixes",
        "bullets": [
            "List 2-3 known bugs or limitations",
            "Explain how to reproduce and current mitigation/fix plan"
        ],
        "notes": "Keep it honest and short; emphasize what you already fixed."
    },
    {
        "title": "Deployment & Next Steps",
        "bullets": [
            "Release steps: signed AAB, Play Console upload",
            "Planned improvements: performance, features, tests"
        ],
        "notes": "Wrap up with realistic next steps and timeline."
    },
    {
        "title": "Git Workflow & Demo Commands",
        "bullets": [
            "Branch: git checkout -b presentation",
            "Build: ./gradlew assembleDebug",
            "Run tests: ./gradlew test"
        ],
        "notes": "Quick commands slide so you can reference them during Q&A if needed."
    },
    {
        "title": "Q & A",
        "bullets": [
            "Thank you — questions welcome"
        ],
        "notes": "Invite questions and be ready to show code or logs for any ask."
    }
]

prs = Presentation()

for i, s in enumerate(slides):
    # Title slide
    if i == 0:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = s.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(52)
        subtitle.text = s.get('subtitle', '')
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    else:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        title.text = s.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(40)

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        bullets = s.get('bullets', [])
        for j, b in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
                p.text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            p.level = 0
            p.font.size = Pt(20)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Add speaker notes
    notes = s.get('notes')
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

out_dir = os.path.join(os.path.dirname(__file__), '')
out_path = os.path.join(out_dir, 'MagdyDiner_Presentation.pptx')
prs.save(out_path)
print(f'Saved {out_path}')

